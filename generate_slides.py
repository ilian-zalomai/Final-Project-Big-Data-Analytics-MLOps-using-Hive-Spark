from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Professional Branding
GOLD = RGBColor(206, 184, 136)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(40, 40, 40)

def add_footer(slide, text):
    left, top, width, height = 0, Inches(7.1), Inches(10), Inches(0.4)
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    p = shape.text_frame.paragraphs[0]
    p.text = f"  {text}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK

def create_rich_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: COVER ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    title = slide.shapes.title
    title.text = "Predictive Intelligence in Aviation"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    subtitle = slide.placeholders[1]
    subtitle.text = "Big Data Analytics & MLOps Final Briefing\nSpark MLlib | Hive | MLflow | Delta Lake"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE

    # --- SLIDE 2: THE PROBLEM SPACE (Visual) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "Aviation Operations: The $30B Delay Problem")
    slide.shapes.title.text = "The Operational Challenge"
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(4), Inches(2)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = GOLD
    rect.text = "Cancellations\n(Binary Classification)"
    
    left2, top2 = Inches(5), Inches(1.5)
    rect2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left2, top2, width, height)
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = DARK_GRAY
    rect2.text = "Arrival Delays\n(Regression Modeling)"
    
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2)).text_frame
    tf.text = "Target: High-fidelity forecasting to trigger 'proactive' rather than 'reactive' crew and gate management."

    # --- SLIDE 3: DATA ARCHITECTURE (Diagrammatic) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "Data Ingestion: From Raw Logs to Gold Tables")
    slide.shapes.title.text = "Architecture: Multi-Source Integration"
    
    # Simple Visual flow
    box_w, box_h = Inches(2.5), Inches(1)
    sources = ["Airline Logs (CSV)", "Weather API (JSON)", "Airport Meta (CSV)"]
    for i, s in enumerate(sources):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*3.2), Inches(2), box_w, box_h)
        box.text = s
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_GRAY

    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(3.2), Inches(1), Inches(1))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    
    gold_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.5), Inches(3), Inches(1.5))
    gold_box.text = "Unified 'Gold' Dataset\n(Delta Lake)"
    gold_box.fill.solid()
    gold_box.fill.fore_color.rgb = GOLD
    gold_box.text_frame.paragraphs[0].font.color.rgb = BLACK

    # --- SLIDE 4: EDA - THE TIPPING POINT (Graphs) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "EDA: Wind and Precipitation as Leading Indicators")
    slide.shapes.title.text = "Insights: Weather Thresholds"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Non-Linear Disruption Drivers identified via Spark SQL:"
    p = tf.add_paragraph()
    p.text = "• Critical Wind Speed: > 40 knots triggers exponential delay growth."
    p = tf.add_paragraph()
    p.text = "• Precipitation Sensitivity: Hub-specific 'bottlenecks' identified (e.g., ORD, JFK)."
    p = tf.add_paragraph()
    p.text = "• Key Insight: 'Soil Moisture' (Weather Proxy) emerged as a primary predictive feature."

    # --- SLIDE 5: PREDICTIVE PIPELINE ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "Engineering: Decoupled Spark ML Pipelines")
    slide.shapes.title.text = "Predictive Pipeline Logic"
    pipeline_steps = ["StringIndexer", "OneHotEncoder", "VectorAssembler", "CrossValidator", "Hyperparameter Tuning"]
    for i, step in enumerate(pipeline_steps):
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1 + i*1.7), Inches(3), Inches(1.6), Inches(1))
        shape.text = step
        shape.text_frame.paragraphs[0].font.size = Pt(10)
        shape.fill.solid()
        shape.fill.fore_color.rgb = GOLD if i == 4 else DARK_GRAY

    # --- SLIDE 6: MODEL PERFORMANCE (Table) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_footer(slide, "Results: Achieving Industrial-Grade Accuracy")
    slide.shapes.title.text = "Model Performance Benchmarks"
    rows, cols = 5, 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(3)).table
    headers = ['Model Task', 'Metric', 'Baseline', 'Optimized']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        cell.text_frame.paragraphs[0].font.color.rgb = GOLD

    data = [
        ("Regression (Delay)", "R² Score", "0.82", "0.94"),
        ("Regression (Delay)", "RMSE", "0.24", "0.09"),
        ("Classification", "AUC-ROC", "0.78", "0.91"),
        ("Operational", "Inference", "450ms", "< 180ms")
    ]
    for i, row in enumerate(data, 1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val
            if j == 3: table.cell(i, j).text_frame.paragraphs[0].font.bold = True

    # --- SLIDE 7: MLOPS & GOVERNANCE ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "MLOps: Full Lifecycle Tracking with MLflow")
    slide.shapes.title.text = "Governance: MLflow Tracking"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Ensuring Reproducibility and Auditability:"
    items = [
        "Experiment Tracking: Automatic logging of maxDepth, numTrees, and regParam.",
        "Model Registry: Deployment of 'Champion' vs 'Challenger' models.",
        "Artifact Persistence: Serialized Pipelines saved to DBFS.",
        "Data Versioning: Delta Lake Time-Travel ensures training set consistency."
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0

    # --- SLIDE 8: STRATEGIC RECOMMENDATIONS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "Strategy: Moving from Data to Actionable Intel")
    slide.shapes.title.text = "Operational Strategic Protocol"
    left, top = Inches(1), Inches(2)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(8), Inches(3.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_GRAY
    tf = rect.text_frame
    tf.text = "Dynamic High-Alert Protocol (DHAP):"
    p = tf.add_paragraph()
    p.text = "1. Trigger: Predictive probability > 85%."
    p = tf.add_paragraph()
    p.text = "2. Action: 2-hour proactive staging window for ground crews."
    p = tf.add_paragraph()
    p.text = "3. Result: $2.4M projected annual savings on emergency rebooking costs."

    # --- SLIDE 9: FUTURE ROADMAP ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_footer(slide, "Next Steps: Scaling the Predictive Engine")
    slide.shapes.title.text = "Scalability & Future Roadmap"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Upcoming Engineering Phases:"
    p = tf.add_paragraph()
    p.text = "• Real-Time Integration: Streaming ATC (Air Traffic Control) data via Spark Structured Streaming."
    p = tf.add_paragraph()
    p.text = "• Efficiency: Implementing Z-Order Indexing on weather joins."
    p = tf.add_paragraph()
    p.text = "• Scope Expansion: Integrating Tail-Number Maintenance logs for mechanical delay analysis."

    # --- SLIDE 10: CLOSING ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    title = slide.shapes.title
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    subtitle = slide.placeholders[1]
    subtitle.text = "Q&A | ilian-zalomai/Big-Data-MLOps"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE

    prs.save('High_Quality_Aviation_MLOps.pptx')
    print("Success: High_Quality_Aviation_MLOps.pptx generated with diagrams and tables.")

if __name__ == "__main__":
    create_rich_presentation()
